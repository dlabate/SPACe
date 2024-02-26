import numpy as np
from pathlib import WindowsPath
from skimage import  io
np.set_printoptions(linewidth=500)
path = WindowsPath(r'E:\tmp\pankaj\20220912-CP-Bolt_20220912_153142\results')
list_pngs = list(path.glob("*.png"))
import pptx
import pptx.util
import glob
import scipy.misc

OUTPUT_TAG = "MY_TAG"

# new
prs = pptx.Presentation()
# open
# prs_exists = pptx.Presentation("some_presentation.pptx")

# default slide width
prs.slide_width = 9144000
# slide height @ 4:3
#prs.slide_height = 6858000
# slide height @ 16:9
# prs.slide_height = 5143500

# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
# blank slide
#slide = prs.slides.add_slide(prs.slide_layouts[6])

# set title
title = slide.shapes.title
title.text = OUTPUT_TAG

pic_left  = int(prs.slide_width * 0.05)
pic_top   = int(prs.slide_height * 0.05)
pic_width = int(prs.slide_width * 0.95)

for g in list_pngs:
    # print(g)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
    # p = tb.textframe.add_paragraph()
    # p.text = g
    # p.font.size = pptx.util.Pt(14)

    img = io.imread(str(g))
    pic_height = int(pic_width *(img.shape[0] / img.shape[1]))
    # pic   = slide.shapes.add_picture(str(g), pic_left, pic_top)
    pic   = slide.shapes.add_picture(str(g), pic_left, pic_top, pic_width, pic_height)


prs.save("%s.pptx" % OUTPUT_TAG)